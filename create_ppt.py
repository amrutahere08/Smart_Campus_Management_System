from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define some branding colors
    primary_color = RGBColor(0, 51, 102)     # Dark Blue
    secondary_color = RGBColor(0, 153, 204)  # Light Blue
    accent_color = RGBColor(255, 102, 0)     # Orange/Gold

    def set_font(run, size, bold=False, color=None):
        run.font.size = Pt(size)
        run.font.name = 'Arial'
        run.font.bold = bold
        if color:
            run.font.color.rgb = color

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "AI-Powered University Assistant and\nSmart Campus Monitoring System"
    subtitle.text = "Project Team:\nAbhilash G [24PG00068]\nAmruta Girish Hegde\nVarshitha R\n\nProject Guide: Mr. Deepak B"

    # Style Title
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            set_font(run, 32, bold=True, color=primary_color)
    
    # Style Subtitle
    for paragraph in subtitle.text_frame.paragraphs:
        for run in paragraph.runs:
            set_font(run, 18, color=RGBColor(80, 80, 80))


    # Helper to add content slides
    def add_content_slide(title_text, content_items):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        # Style Title
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = primary_color
        
        # Body
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # Clear default empty paragraph

        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.space_after = Pt(14)
            p.level = 0

    # Slide 2: Introduction
    intro_points = [
        "The Smart Campus project integrates Artificial Intelligence (AI) and Internet of Things (IoT) to enhance the university experience.",
        "It aims to bridge the gap between physical campus administration and digital convenience.",
        "Key focus areas include efficient visitor management, automated student assistance, and real-time monitoring."
    ]
    add_content_slide("Introduction", intro_points)

    # Slide 3: Problem Statement
    problem_points = [
        "Manual visitor entry logs are time-consuming, prone to errors, and insecure.",
        "Students frequent administrative blocks for basic queries (schedules, faculty location), causing congestion.",
        "Faculty members find it difficult to track appointments and student interactions efficiently.",
        "Lack of centralized real-time insights for campus administrators."
    ]
    add_content_slide("Problem Statement", problem_points)

    # Slide 4: Proposed Solution
    solution_points = [
        "Visitor Kiosk: Automated entry using Face Recognition and Emotion Detection.",
        "AI Assistant: A smart chatbot (Voice & Text) to answer queries about faculties, courses, and facilities.",
        "Admin Dashboard: comprehensive control panel for User Management and Security Monitoring.",
        "Faculty Portal: Dedicated interface for faculty to manage availability and view visitor logs."
    ]
    add_content_slide("Proposed Solution", solution_points)

    # Slide 5: System Architecture
    architecture_points = [
        "Frontend: HTML5, CSS3, JavaScript (Responsive & Modern UI)",
        "Backend: Python (Flask Framework)",
        "Database: SQLite (scalable to PostgreSQL/MySQL)",
        "AI Models: DeepFace (Face Recognition), Custom NLP Models (Chatbot)",
        "Hardware Integration: Camera feeds for monitoring and kiosk interaction."
    ]
    add_content_slide("System Architecture", architecture_points)

    # Slide 6: Key Features
    features_points = [
        "Face Identification: Seamless and secure entry for registered users.",
        "Emotion Analysis: Detects visitor mood to personalize greetings.",
        "Smart Navigation: Provides directions to library, labs, and faculty rooms.",
        "Role-Based Access Control: Specialized views for Admin, Faculty, Security, and Students."
    ]
    add_content_slide("Key Features", features_points)

    # Slide 7: Individual Contributions
    # Note: These are placeholders based on typical role distribution.
    contributions_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(contributions_slide_layout)
    title = slide.shapes.title
    title.text = "Individual Contributions"
    title.text_frame.paragraphs[0].runs[0].font.color.rgb = primary_color
    
    body = slide.placeholders[1]
    tf = body.text_frame
    
    # Abhilash
    p = tf.text = "Abhilash G [24PG00068]"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Backend Architecture (Flask)\n• Database Schema Design\n• System Integration & Security Modules"
    p.level = 1
    
    # Amruta
    p = tf.add_paragraph()
    p.text = "Amruta Girish Hegde"
    p.font.bold = True
    p.font.size = Pt(20)
    p.space_before = Pt(20)

    p = tf.add_paragraph()
    p.text = "• Frontend UI/UX Design\n• Visitor Kiosk Interface & Animation\n• Emotion Detection Model Integration"
    p.level = 1
    
    # Varshitha
    p = tf.add_paragraph()
    p.text = "Varshitha R"
    p.font.bold = True
    p.font.size = Pt(20)
    p.space_before = Pt(20)

    p = tf.add_paragraph()
    p.text = "• AI Assistant Logic & NLP Training\n• Testing & Validation\n• Documentation & Content"
    p.level = 1

    # Slide 8: Technology Stack (Visuals Placeholder)
    tech_points = [
        "Python", "Flask", "OpenCV", "DeepFace", "HTML/CSS/JS", "SQLite"
    ]
    add_content_slide("Technology Stack", tech_points)

    # Slide 9: Conclusion
    conclusion_points = [
        "The project successfully demonstrates a scalable Smart Campus solution.",
        "It reduces manual workload for security and admin staff.",
        "Enhances the technological appeal of the university.",
        "Ready for future expansions like mobile app integration and predictive analytics."
    ]
    add_content_slide("Conclusion", conclusion_points)

    # Slide 10: Thank You
    slide_layout = prs.slide_layouts[0] # Title Slide format
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Thank You!"
    title.top = Inches(2.5)
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Any Questions?"

    filename = "Smart_Campus_Project_Presentation.pptx"
    prs.save(filename)
    print(f"Presentation saved to {filename}")

if __name__ == "__main__":
    create_presentation()

#!/usr/bin/env python3
"""
Create PowerPoint Presentation for Smart Campus Project
Documents all AI/ML models used, their locations, and purposes
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Smart Campus Project"
    subtitle.text = "AI/ML Models Documentation\n\nA Comprehensive Overview of Machine Learning Models\nUsed in the Smart Campus System"
    
    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def add_overview_slide(prs):
    """Add project overview slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Overview"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    points = [
        ("Smart Campus System", "AI-powered campus management platform built with Flask"),
        ("Purpose", "Automate attendance, visitor management, and campus assistance"),
        ("Technology Stack", "Python, Flask, SQLite, Multiple AI/ML Models"),
        ("Key Features", "Face Recognition, Emotion Detection, AI Chatbot, Voice Interaction, Attention Monitoring"),
        ("Architecture", "Modular Monolith with Service Layer Pattern")
    ]
    
    for heading, text in points:
        p = content.add_paragraph()
        p.text = f"{heading}: {text}"
        p.level = 0
        p.font.size = Pt(16)
        p.space_after = Pt(12)

def add_models_summary_slide(prs):
    """Add models summary slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "AI/ML Models Used - Summary"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    models = [
        "1. Google Gemini 2.5 Flash - AI Chatbot",
        "2. face_recognition (dlib) - Face Recognition",
        "3. DeepFace - Emotion Detection",
        "4. MediaPipe Face Landmarker - Attention Monitoring",
        "5. Google Speech Recognition - Voice-to-Text",
        "6. gTTS (Google Text-to-Speech) - Text-to-Voice"
    ]
    
    for model in models:
        p = content.add_paragraph()
        p.text = model
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(14)

def add_model_detail_slide(prs, model_name, details):
    """Add detailed slide for a specific model"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = model_name
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    for key, value in details.items():
        # Add heading
        p = content.add_paragraph()
        p.text = key
        p.level = 0
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.space_after = Pt(6)
        
        # Add content
        if isinstance(value, list):
            for item in value:
                p = content.add_paragraph()
                p.text = item
                p.level = 1
                p.font.size = Pt(14)
                p.space_after = Pt(4)
        else:
            p = content.add_paragraph()
            p.text = value
            p.level = 1
            p.font.size = Pt(14)
            p.space_after = Pt(8)

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs)
    
    # Slide 2: Overview
    add_overview_slide(prs)
    
    # Slide 3: Models Summary
    add_models_summary_slide(prs)
    
    # Slide 4: Google Gemini
    add_model_detail_slide(prs, "Model 1: Google Gemini 2.5 Flash", {
        "Purpose": "AI-powered chatbot for campus assistance and information",
        "Location": "services/chatbot.py (ChatbotService class)",
        "Model Type": "Large Language Model (LLM) - Generative AI",
        "Key Features": [
            "Natural language understanding and generation",
            "Context-aware responses about campus events, faculty, departments",
            "Real-time availability checking for students and faculty",
            "Location information for faculty cabins and offices",
            "Integration with campus database for accurate information"
        ],
        "Why This Model": [
            "Latest Gemini 2.5 Flash for fast, accurate responses",
            "Excellent at understanding context and providing relevant information",
            "Free tier available for development and testing",
            "Strong performance on question-answering tasks"
        ],
        "API Used": "google.generativeai library with Gemini API"
    })
    
    # Slide 5: Face Recognition
    add_model_detail_slide(prs, "Model 2: face_recognition (dlib)", {
        "Purpose": "Automatic face detection, enrollment, and recognition for attendance",
        "Location": "services/face_recognition.py (FaceRecognitionService class)",
        "Model Type": "Deep Learning Face Recognition (dlib's ResNet-based model)",
        "Key Features": [
            "Face detection in images and video streams",
            "128-dimensional face encoding generation",
            "Automatic face matching with enrolled faces",
            "Attendance tracking with IN/OUT detection",
            "Duplicate entry prevention (5-minute threshold)"
        ],
        "Why This Model": [
            "Industry-standard accuracy (99.38% on LFW benchmark)",
            "Fast inference suitable for real-time applications",
            "Works well with varying lighting conditions",
            "Robust to pose variations and facial expressions",
            "Open-source and well-documented"
        ],
        "Technical Details": "Uses dlib's CNN-based face detector and ResNet-based face encoder"
    })
    
    # Slide 6: DeepFace
    add_model_detail_slide(prs, "Model 3: DeepFace", {
        "Purpose": "Emotion detection and demographic analysis for personalized greetings",
        "Location": "services/emotion_detection.py (EmotionDetectionService class)",
        "Model Type": "Multi-task Deep Learning Model (Emotion, Age, Gender)",
        "Key Features": [
            "Detects 7 emotions: happy, sad, angry, surprise, fear, disgust, neutral",
            "Age estimation",
            "Gender classification",
            "Confidence scoring for predictions",
            "Personalized greeting generation based on detected emotion"
        ],
        "Why This Model": [
            "Comprehensive facial analysis in single inference",
            "Pre-trained on large facial expression datasets",
            "High accuracy for emotion recognition",
            "Easy integration with OpenCV",
            "Supports multiple backend models (VGG-Face, Facenet, etc.)"
        ],
        "Use Cases": [
            "Visitor kiosk personalized greetings",
            "Emotion tracking linked to attendance records",
            "User experience enhancement"
        ]
    })
    
    # Slide 7: MediaPipe
    add_model_detail_slide(prs, "Model 4: MediaPipe Face Landmarker", {
        "Purpose": "Real-time student attention monitoring in classrooms",
        "Location": "services/attention_monitoring.py (AttentionMonitoringService class)",
        "Model Type": "Facial Landmark Detection and Head Pose Estimation",
        "Key Features": [
            "Detects up to 30 faces simultaneously",
            "468 facial landmarks per face",
            "Head pose estimation (pitch, yaw, roll angles)",
            "Determines if students are looking at the board or distracted",
            "Alert system when >10 students are distracted"
        ],
        "Why This Model": [
            "Google's production-ready solution for face mesh",
            "Extremely fast - suitable for real-time video processing",
            "Accurate head pose estimation",
            "Handles multiple faces efficiently",
            "Cross-platform support (CPU/GPU)"
        ],
        "Technical Approach": [
            "Uses 6 key facial landmarks for pose estimation",
            "PnP (Perspective-n-Point) algorithm for 3D pose calculation",
            "Thresholds: Pitch >30° (looking down), Yaw >30° (looking away)"
        ]
    })
    
    # Slide 8: Google Speech Recognition
    add_model_detail_slide(prs, "Model 5: Google Speech Recognition", {
        "Purpose": "Convert voice input to text for voice-based chatbot interaction",
        "Location": "services/voice_chat.py (VoiceChatService class)",
        "Model Type": "Automatic Speech Recognition (ASR)",
        "Key Features": [
            "Converts audio files to text",
            "Supports multiple languages",
            "Handles various audio formats (WAV, MP3)",
            "Integration with chatbot for voice queries",
            "Real-time speech-to-text conversion"
        ],
        "Why This Model": [
            "High accuracy speech recognition",
            "Free tier available via SpeechRecognition library",
            "Robust to background noise",
            "Supports natural language input",
            "Easy integration with Python"
        ],
        "Library Used": "speech_recognition (SpeechRecognition) with Google Web Speech API",
        "Workflow": "Audio → Speech Recognition → Text → Chatbot → Response"
    })
    
    # Slide 9: gTTS
    add_model_detail_slide(prs, "Model 6: gTTS (Google Text-to-Speech)", {
        "Purpose": "Convert chatbot text responses to natural-sounding speech",
        "Location": "services/voice_chat.py and services/voice_guidance.py",
        "Model Type": "Text-to-Speech (TTS) Synthesis",
        "Key Features": [
            "Converts text responses to audio",
            "Natural-sounding voice output",
            "Multiple language support",
            "Adjustable speech rate",
            "Voice guidance for visitor kiosk"
        ],
        "Why This Model": [
            "High-quality, natural-sounding voices",
            "Free and easy to use",
            "Supports 100+ languages",
            "Reliable and fast",
            "Perfect for accessibility features"
        ],
        "Use Cases": [
            "Voice chatbot responses",
            "Visitor kiosk voice prompts",
            "Accessibility for visually impaired users",
            "Interactive voice guidance system"
        ],
        "Library Used": "gTTS (Google Text-to-Speech) Python library"
    })
    
    # Slide 10: Architecture Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "System Architecture"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    arch_points = [
        ("Framework", "Flask (Python Web Framework)"),
        ("Database", "SQLite with SQLAlchemy ORM"),
        ("Architecture Pattern", "Modular Monolith with Service Layer"),
        ("AI Services Layer", "Encapsulated in services/ directory"),
        ("Face Data Storage", "Hybrid: Images in filesystem, encodings in database"),
        ("Processing", "Local AI processing for privacy and low latency"),
        ("API Endpoints", "RESTful APIs in routes/api/ for each service")
    ]
    
    for heading, text in arch_points:
        p = content.add_paragraph()
        p.text = f"{heading}: {text}"
        p.level = 0
        p.font.size = Pt(15)
        p.space_after = Pt(10)
    
    # Slide 11: Model Integration Flow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Model Integration Flow"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    p = content.add_paragraph()
    p.text = "Face Recognition Flow"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    flow1 = [
        "1. Camera captures image → OpenCV",
        "2. Face detection → face_recognition library",
        "3. Generate 128-d encoding → dlib model",
        "4. Compare with enrolled faces → Database query",
        "5. Match found → Create attendance record",
        "6. Emotion analysis → DeepFace",
        "7. Store emotion data → Link to attendance"
    ]
    
    for step in flow1:
        p = content.add_paragraph()
        p.text = step
        p.level = 1
        p.font.size = Pt(14)
        p.space_after = Pt(6)
    
    # Slide 12: Voice Interaction Flow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Voice Interaction Flow"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    p = content.add_paragraph()
    p.text = "Complete Voice Chat Pipeline"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    flow2 = [
        "1. User speaks → Microphone captures audio",
        "2. Audio processing → speech_recognition library",
        "3. Speech-to-Text → Google Speech Recognition API",
        "4. Text query → ChatbotService",
        "5. Context building → Database queries (events, faculty, etc.)",
        "6. AI processing → Google Gemini 2.5 Flash",
        "7. Text response → Generated answer",
        "8. Text-to-Speech → gTTS",
        "9. Audio playback → User hears response"
    ]
    
    for step in flow2:
        p = content.add_paragraph()
        p.text = step
        p.level = 1
        p.font.size = Pt(13)
        p.space_after = Pt(5)
    
    # Slide 13: Key Benefits
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Benefits of AI/ML Integration"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    benefits = [
        ("Automation", "Automated attendance tracking eliminates manual processes"),
        ("Accuracy", "99%+ face recognition accuracy with dlib-based model"),
        ("User Experience", "Personalized greetings based on emotion detection"),
        ("Accessibility", "Voice interaction for hands-free campus assistance"),
        ("Real-time Monitoring", "Live attention tracking for classroom engagement"),
        ("Privacy", "Local AI processing keeps sensitive data on-premise"),
        ("Scalability", "Modular architecture allows easy addition of new AI features"),
        ("Intelligence", "Context-aware chatbot provides accurate campus information")
    ]
    
    for heading, text in benefits:
        p = content.add_paragraph()
        p.text = f"{heading}: {text}"
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(10)
    
    # Slide 14: Technical Specifications
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technical Specifications"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    specs = [
        ("Programming Language", "Python 3.10+"),
        ("Web Framework", "Flask with Blueprints"),
        ("Database", "SQLite (easily switchable to PostgreSQL/MySQL)"),
        ("ORM", "SQLAlchemy"),
        ("Authentication", "Flask-Login with bcrypt password hashing"),
        ("Computer Vision", "OpenCV, face_recognition, DeepFace, MediaPipe"),
        ("AI/ML", "Google Gemini API, dlib, TensorFlow (via DeepFace)"),
        ("Voice Processing", "SpeechRecognition, gTTS"),
        ("Frontend", "Jinja2 templates, Vanilla JavaScript, CSS3"),
        ("Deployment", "WSGI-compatible (Gunicorn, uWSGI)")
    ]
    
    for heading, text in specs:
        p = content.add_paragraph()
        p.text = f"{heading}: {text}"
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(8)
    
    # Slide 15: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion"
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    conclusion_points = [
        "Smart Campus leverages 6 state-of-the-art AI/ML models",
        "Each model serves a specific, well-defined purpose",
        "Integration creates a seamless, intelligent campus experience",
        "Local processing ensures privacy and low latency",
        "Modular architecture allows for future enhancements",
        "Production-ready with proven, industry-standard models",
        "Comprehensive solution covering face recognition, emotion detection, voice interaction, and attention monitoring"
    ]
    
    for point in conclusion_points:
        p = content.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(16)
        p.space_after = Pt(12)
    
    # Save presentation
    output_file = "Smart_Campus_AI_Models_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Presentation created successfully: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    create_presentation()
